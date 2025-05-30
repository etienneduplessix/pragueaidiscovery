from fastapi import FastAPI, File, UploadFile, HTTPException, Request, Body, Form, Depends
from fastapi.templating import Jinja2Templates
from fastapi.responses import HTMLResponse, JSONResponse, StreamingResponse, FileResponse
from fastapi.middleware.cors import CORSMiddleware
from sqlalchemy.orm import Session
from sqlalchemy import text, inspect
from typing import List, Dict, Optional
from pydantic import BaseModel
import io
import os
import tempfile
import re
import uuid
import boto3
from botocore.client import Config
from datetime import datetime

# Image processing
from PIL import Image
import pytesseract
from pdf2image import convert_from_path

# Data processing
import pandas as pd
import numpy as np

# PowerPoint
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor

# Import database connection
from database import engine, SessionLocal, Base

# Initialize FastAPI app
app = FastAPI(title="Document Processing API")

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Jinja2 template engine setup
templates = Jinja2Templates(directory="templates")

# Model definitions
class PptxReportPayload(BaseModel):
    table_name: str

class BufferData(BaseModel):
    type: str
    data: List[int]

class FileBody(BaseModel):
    body: BufferData

# MinIO configuration - Consider moving these to environment variables
MINIO_CONFIG = {
    "endpoint": "minio:9000",
    "access_key": "HWch5ike5Vz3flhdVjBr",
    "secret_key": "w4TP1WkuB0DqbyXgpE6YNzOjoKSEJa6HCFkLC7H3",
    "bucket": "uploads"
}

# Initialize MinIO S3 client
s3_client = boto3.client(
    "s3",
    endpoint_url=f"http://{MINIO_CONFIG['endpoint']}",
    aws_access_key_id=MINIO_CONFIG['access_key'],
    aws_secret_access_key=MINIO_CONFIG['secret_key'],
    config=Config(signature_version="s3v4"),
    region_name="us-east-1"  # Dummy region
)

# Database dependency
def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

# Helper functions
def process_image_ocr(image):
    """Extract text from an image using OCR"""
    return pytesseract.image_to_string(image)

def process_pdf_ocr(pdf_content):
    """Extract text from a PDF using OCR"""
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
        temp_file.write(pdf_content)
        temp_path = temp_file.name
    
    extracted_text = ""
    try:
        images = convert_from_path(temp_path)
        for image in images:
            extracted_text += pytesseract.image_to_string(image) + "\n"
    finally:
        os.unlink(temp_path)
    
    return extracted_text

def get_content_type(filename):
    """Determine content type based on file extension"""
    extension = filename.split('.')[-1].lower()
    content_types = {
        'pdf': 'application/pdf', 
        'png': 'image/png',
        'jpg': 'image/jpg',
        'jpeg': 'image/jpeg',
        'gif': 'image/gif',
        'csv': 'text/csv',
    }
    return content_types.get(extension, 'application/octet-stream')

def validate_table_name(table_name):
    """Validate a table name to prevent SQL injection"""
    if not re.fullmatch(r"csv_[a-zA-Z0-9_]+", table_name):
        raise HTTPException(status_code=400, detail="Invalid table name")
    return table_name

# Routes
@app.get("/", response_class=HTMLResponse)
def read_root(request: Request):
    """Serve the main page"""
    return templates.TemplateResponse("index.html", {"request": request})

@app.post("/ocr1/")
async def extract_text_from_minio(filename: str):
    """Extract text from a file stored in MinIO"""
    try:
        # Download file from MinIO
        file_object = s3_client.get_object(Bucket=MINIO_CONFIG['bucket'], Key=filename)
        content = file_object['Body'].read()
        
        # Determine file type
        if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
            image = Image.open(io.BytesIO(content))
            extracted_text = process_image_ocr(image)
            return {"extracted_text": extracted_text}
        
        elif filename.lower().endswith('.pdf'):
            extracted_text = process_pdf_ocr(content)
            return {"extracted_text": extracted_text}
        
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type for file: {filename}")
        
    except s3_client.exceptions.NoSuchKey:
        raise HTTPException(status_code=404, detail=f"File not found in storage: {filename}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

@app.post("/ocr2/")
async def extract_text(
    file: UploadFile = File(None),
    file_body: FileBody = Body(None)
):
    """Extract text from a file uploaded via multipart/form-data or JSON payload"""
    try:
        # Handle binary upload via multipart/form-data
        if file:
            content = await file.read()
            file_extension = file.filename.split('.')[-1].lower()
        
        # Handle JSON payload with body.data
        elif file_body:
            content = bytes(file_body.body.data)
            # Simple file type detection based on magic bytes
            if content[:4] == b'\x89PNG':
                file_extension = 'png'
            elif content[:3] == b'GIF':
                file_extension = 'gif'
            elif content[:2] == b'\xff\xd8':
                file_extension = 'jpg'
            elif content[:4] == b'%PDF':
                file_extension = 'pdf'
            else:
                file_extension = 'unknown'
        else:
            raise HTTPException(status_code=400, detail="No file provided")
            
        # Process files based on type
        if file_extension.lower() in ['png', 'jpg', 'jpeg', 'gif']:
            image = Image.open(io.BytesIO(content))
            extracted_text = process_image_ocr(image)
            return {"extracted_text": extracted_text}
            
        elif file_extension.lower() == 'pdf':
            extracted_text = process_pdf_ocr(content)
            return {"extracted_text": extracted_text}
            
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {file_extension}")

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

@app.get("/api/tables")
def list_csv_tables():
    """List all CSV tables in the database"""
    inspector = inspect(engine)
    all_tables = inspector.get_table_names()
    csv_tables = [t for t in all_tables if t.startswith("csv_")]
    return {"tables": csv_tables}

@app.get("/api/table/{table_name}")
def get_table_data(table_name: str):
    """Retrieve data from a specific CSV table"""
    table_name = validate_table_name(table_name)
    
    try:
        query = f"SELECT * FROM {table_name} LIMIT 100"
        with engine.connect() as conn:
            result = conn.execute(text(query))
            rows = [dict(row._mapping) for row in result]
        return {"data": rows}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Database error: {str(e)}")

@app.get("/download/{filename}")
async def download_file(filename: str):
    """Get file metadata and binary content"""
    try:
        file_object = s3_client.get_object(Bucket=MINIO_CONFIG['bucket'], Key=filename)
        file_data = file_object['Body'].read()

        return {
            "filename": filename,
            "content_size": len(file_data),
            "content": list(file_data)  # Returns binary data as a list of bytes
        }
    except s3_client.exceptions.NoSuchKey:
        raise HTTPException(status_code=404, detail="File not found in storage")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error accessing storage: {str(e)}")

@app.get("/download-file/{filename}")
async def download_file_stream(filename: str):
    """Download a file as a streaming response"""
    try:
        file_object = s3_client.get_object(Bucket=MINIO_CONFIG['bucket'], Key=filename)
        content_type = get_content_type(filename)
        
        return StreamingResponse(
            file_object['Body'],
            media_type=content_type,
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    except s3_client.exceptions.NoSuchKey:
        raise HTTPException(status_code=404, detail="File not found in storage")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error accessing storage: {str(e)}")

@app.post("/upload/")
async def upload_file(
    file: UploadFile = File(...), 
    filetype: str = Form(None), 
    overwrite: bool = Form(False)
):
    """Upload a file to MinIO with optional processing"""
    try:
        # Check if file already exists
        try:
            s3_client.head_object(Bucket=MINIO_CONFIG['bucket'], Key=file.filename)
            if not overwrite:
                return {
                    "filename": file.filename,
                    "status": "already_exists",
                    "message": "File already exists in storage. Set 'overwrite=true' to replace it.",
                    "path": f"/download-file/{file.filename}"
                }
        except s3_client.exceptions.ClientError as e:
            if e.response['Error']['Code'] != '404':
                raise
        
        # Read file content
        file_content = await file.read()
        
        # Determine content type
        content_type = None
        if filetype:
            content_type = get_content_type(filetype)
        
        # Upload to MinIO
        extra_args = {'ContentType': content_type} if content_type else {}
        s3_client.put_object(
            Bucket=MINIO_CONFIG['bucket'],
            Key=file.filename,
            Body=file_content,
            **extra_args
        )
        
        result = {
            "filename": file.filename,
            "size": len(file_content),
            "status": "uploaded_successfully",
            "path": f"/download-file/{file.filename}"
        }
        
        # Process specific file types
        if filetype and filetype.lower() in ['jpg', 'jpeg', 'png', 'gif', 'pdf']:
            # OCR processing
            extracted_text = ""
            if filetype.lower() in ['jpg', 'jpeg', 'png', 'gif']:
                image = Image.open(io.BytesIO(file_content))
                extracted_text = process_image_ocr(image)
            elif filetype.lower() == 'pdf':
                extracted_text = process_pdf_ocr(file_content)
            
            result["extracted_text"] = extracted_text if extracted_text else None
            
        elif filetype and filetype.lower() == 'csv':
            # CSV processing
            try:
                # Create a StringIO object from the file content
                csv_io = io.StringIO(file_content.decode('utf-8'))
                
                # Read the CSV file
                df = pd.read_csv(csv_io)
                
                # Add CSV information to result
                result["csv_preview"] = df.head(5).to_dict('records')
                result["rows"] = len(df)
                result["columns"] = df.columns.tolist()
                
                # Uncomment to save to database if needed
                # table_name = f"csv_{os.path.splitext(file.filename)[0].replace('-', '_').replace(' ', '_')}"
                # df.to_sql(table_name, engine, index=False, if_exists='replace')
                # result["table_name"] = table_name
                
            except Exception as csv_error:
                result["csv_error"] = str(csv_error)
        
        return result
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error uploading file: {str(e)}")

@app.post("/generate-report-pptx/")
async def generate_report_pptx(payload: PptxReportPayload):
    """Generate a PowerPoint report from a database table"""
    table_name = validate_table_name(payload.table_name)
    
    try:
        # Load data from database
        query = f"SELECT * FROM {table_name}"
        with engine.connect() as conn:
            df = pd.read_sql(query, conn)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Database error: {str(e)}")
    
    try:
        # Data processing
        # Clean and convert the total_price column to numeric
        df['total_price'] = pd.to_numeric(df['total_price'].str.replace(',', ''), errors='coerce')
        df['total_price'] = df['total_price'].fillna(0)
        
        # Process dates
        df['purchase_date'] = pd.to_datetime(df['purchase_date'], errors='coerce')
        df = df.dropna(subset=['purchase_date'])
        
        df['month'] = df['purchase_date'].dt.month
        df['quarter'] = df['purchase_date'].dt.quarter
        df['year'] = df['purchase_date'].dt.year

        # Customer repeat analysis
        df = df.sort_values('purchase_date')
        df['first_purchase'] = df.groupby('customer_id')['purchase_date'].transform('min')
        df['is_returning'] = df['purchase_date'] > df['first_purchase']
        repeat_customers = df['customer_id'].value_counts()
        repeat_rate = (repeat_customers > 1).sum() / repeat_customers.count()

        # Create presentation
        prs = Presentation()

        # Utility functions for presentation creation
        def add_section(title):
            slide = prs.slides.add_slide(prs.slide_layouts[2])
            slide.shapes.title.text = title
            return slide

        def add_content(title, content=None):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            if content:
                slide.placeholders[1].text = content
            return slide

        def add_chart(slide, chart_type, categories, values, title=""):
            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series(title, values)
            x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
            chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data).chart
            chart.has_legend = True
            chart.has_title = True
            chart.chart_title.text_frame.text = title

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Sales Analytics Report"
        slide.placeholders[1].text = f"From {table_name} • {datetime.now().strftime('%B %d, %Y')}"

        # 1. TOTAL SALES OVERVIEW
        add_section("Total Sales Overview")
        total_revenue = df['total_price'].sum()
        avg_revenue = df['total_price'].mean()

        slide = add_content("Revenue Stats")
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
        tf = box.text_frame
        for text, size in [
            (f"Total Revenue: ${total_revenue:,.2f}", 24),
            (f"Average Transaction: ${avg_revenue:,.2f}", 18),
        ]:
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(size)

        # Monthly Revenue
        m = df.groupby('month')['total_price'].sum().reset_index()
        add_chart(add_content("Monthly Revenue"), XL_CHART_TYPE.COLUMN_CLUSTERED,
                  [f"Month {x}" for x in m['month']], m['total_price'], "Revenue by Month")

        # Quarterly Revenue
        q = df.groupby('quarter')['total_price'].sum().reset_index()
        add_chart(add_content("Quarterly Revenue"), XL_CHART_TYPE.COLUMN_CLUSTERED,
                  [f"Q{x}" for x in q['quarter']], q['total_price'], "Revenue by Quarter")

        # Growth Rate
        m['growth'] = m['total_price'].pct_change() * 100
        add_chart(add_content("Growth Rate"), XL_CHART_TYPE.LINE,
                  [f"Month {x}" for x in m['month'][1:]], m['growth'][1:], "Month-over-Month Growth (%)")

        # 2. TOP PRODUCTS
        add_section("Top Products/Services")
        prod = df.groupby('product')['total_price'].sum().sort_values(ascending=False).reset_index()
        top_prod = prod.head(10)
        top_prod['contribution'] = top_prod['total_price'] / top_prod['total_price'].sum() * 100

        add_chart(add_content("Top Products"), XL_CHART_TYPE.BAR_CLUSTERED,
                  top_prod['product'], top_prod['total_price'], "Top Products by Revenue")

        add_chart(add_content("Revenue Contribution %"), XL_CHART_TYPE.PIE,
                  top_prod['product'], top_prod['contribution'], "Product Contribution %")

        # 3. CUSTOMER SEGMENTATION
        add_section("Customer Segmentation")

        # Region
        reg = df.groupby('country')['total_price'].sum().sort_values(ascending=False).head(10)
        add_chart(add_content("Revenue by Region"), XL_CHART_TYPE.BAR_CLUSTERED,
                  reg.index.tolist(), reg.tolist(), "Top Regions by Revenue")

        # Top Customers
        cust = df.groupby('customer_name')['total_price'].sum().sort_values(ascending=False).head(10)
        add_chart(add_content("Top Customers"), XL_CHART_TYPE.BAR_CLUSTERED,
                  cust.index.tolist(), cust.tolist(), "Top Customers by Revenue")

        # New vs. Returning
        count = df['is_returning'].value_counts()
        add_chart(add_content("New vs. Returning Customers"), XL_CHART_TYPE.PIE,
                  ['New', 'Returning'],
                  [count.get(False, 0), count.get(True, 0)],
                  "Customer Type Distribution")

        # 4. MARKET TRENDS
        add_section("Market Trends")
        df['period'] = df['year'].astype(str) + '-' + df['month'].astype(str).str.zfill(2)
        ts = df.groupby('period')['total_price'].sum().reset_index()
        add_chart(add_content("Sales Over Time"), XL_CHART_TYPE.LINE,
                  ts['period'], ts['total_price'], "Revenue Trend")

        # Seasonality
        season = df.groupby('month')['total_price'].sum().reset_index()
        add_chart(add_content("Seasonality"), XL_CHART_TYPE.COLUMN_CLUSTERED,
                  [f"Month {m}" for m in season['month']], season['total_price'], "Monthly Seasonality")

        # 5. GEOGRAPHICAL PERFORMANCE
        add_section("Geographical Performance")
        geo = df.groupby('country')['total_price'].sum().sort_values(ascending=False).head(15)
        add_chart(add_content("Geo Performance"), XL_CHART_TYPE.BAR_CLUSTERED,
                  geo.index.tolist(), geo.tolist(), "Revenue by Country")

        # 6. OPPORTUNITIES & ANOMALIES
        add_section("Opportunities & Anomalies")
        under = prod.tail(10)
        add_chart(add_content("Underperforming Products"), XL_CHART_TYPE.BAR_CLUSTERED,
                  under['product'], under['total_price'], "Lowest Revenue Products")

        # SUMMARY SLIDE
        slide = add_content("Summary & Key Metrics")
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3.5))
        tf = box.text_frame
        for label, value in [
            ("Total Revenue", f"${total_revenue:,.2f}"),
            ("Repeat Purchase Rate", f"{repeat_rate:.2%}"),
            ("Unique Products", df['product'].nunique()),
            ("Customers", df['customer_id'].nunique()),
            ("Countries Served", df['country'].nunique()),
        ]:
            p = tf.add_paragraph()
            p.text = f"{label}: {value}"
            p.font.size = Pt(16)

        # Export
        output_path = f"/tmp/{uuid.uuid4()}.pptx"
        prs.save(output_path)
        return FileResponse(
            output_path, 
            filename="sales_report.pptx", 
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Report generation failed: {str(e)}")